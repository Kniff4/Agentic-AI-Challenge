from ibm_watsonx_orchestrate.agent_builder.tools import tool
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
import json
import os
from typing import Any, Dict, List, Optional


def _safe_get(data: Dict[str, Any], key: str, default=None):
    value = data.get(key, default)
    return default if value is None else value


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert #RRGGBB to python-pptx RGBColor."""
    hex_color = hex_color.strip().lstrip("#")
    if len(hex_color) != 6:
        # fallback: Talentia blue
        return RGBColor(36, 19, 95)
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def _style_value(global_style: Dict[str, Any], slide_style: Optional[Dict[str, Any]], key: str, default=None):
    if slide_style and key in slide_style and slide_style[key] is not None:
        return slide_style[key]
    return _safe_get(global_style, key, default)


def _apply_run_style(
    run,
    font_family: Optional[str] = None,
    font_size_pt: Optional[int] = None,
    bold: Optional[bool] = None,
    color_hex: Optional[str] = None,
):
    if font_family:
        run.font.name = font_family
    if font_size_pt:
        run.font.size = Pt(font_size_pt)
    if bold is not None:
        run.font.bold = bold
    if color_hex:
        run.font.color.rgb = _hex_to_rgb(color_hex)


def _apply_paragraph_style(
    paragraph,
    font_family: Optional[str] = None,
    font_size_pt: Optional[int] = None,
    bold: Optional[bool] = None,
    color_hex: Optional[str] = None,
):
    # Ensure there is at least one run
    if not paragraph.runs:
        run = paragraph.add_run()
        run.text = paragraph.text
        paragraph.text = ""
    for run in paragraph.runs:
        _apply_run_style(run, font_family, font_size_pt, bold, color_hex)


def _set_title_text(shape, text: str, global_style: Dict[str, Any], slide_style: Optional[Dict[str, Any]] = None):
    shape.text = text
    if shape.text_frame and shape.text_frame.paragraphs:
        p = shape.text_frame.paragraphs[0]
        _apply_paragraph_style(
            p,
            font_family=_style_value(global_style, slide_style, "title_font_family", "Arial"),
            font_size_pt=_style_value(global_style, slide_style, "title_font_size_pt", 28),
            bold=True,
            color_hex=_style_value(global_style, slide_style, "title_color", _style_value(global_style, slide_style, "primary_color", "#24135F")),
        )


def _set_subtitle_text(shape, text: str, global_style: Dict[str, Any], slide_style: Optional[Dict[str, Any]] = None):
    shape.text = text
    if shape.text_frame and shape.text_frame.paragraphs:
        p = shape.text_frame.paragraphs[0]
        _apply_paragraph_style(
            p,
            font_family=_style_value(global_style, slide_style, "body_font_family", "Arial"),
            font_size_pt=_style_value(global_style, slide_style, "subtitle_font_size_pt", 20),
            bold=False,
            color_hex=_style_value(global_style, slide_style, "body_color", "#24135F"),
        )


def _add_notes(slide, slide_spec: Dict[str, Any]) -> None:
    notes = _safe_get(slide_spec, "notes", "")
    if not notes:
        return
    try:
        slide.notes_slide.notes_text_frame.text = notes
    except Exception:
        pass


def _add_logo(slide, global_style: Dict[str, Any], slide_style: Optional[Dict[str, Any]] = None):
    logo_file = _style_value(global_style, slide_style, "logo_file", "")
    logo_position = _style_value(global_style, slide_style, "logo_position", "top_right")

    if not logo_file or not os.path.exists(logo_file):
        return

    # Conservative default sizing for MVP
    width = Inches(1.2)

    if logo_position == "top_right":
        left = Inches(8.0)
        top = Inches(0.2)
    elif logo_position == "top_left":
        left = Inches(0.3)
        top = Inches(0.2)
    elif logo_position == "bottom_right":
        left = Inches(8.0)
        top = Inches(6.7)
    else:
        left = Inches(8.0)
        top = Inches(0.2)

    try:
        slide.shapes.add_picture(logo_file, left, top, width=width)
    except Exception:
        pass


def _truncate_bullets(bullets: List[Any], max_bullets: int) -> List[Any]:
    return bullets[:max_bullets] if max_bullets and len(bullets) > max_bullets else bullets


def _add_title_slide(prs: Presentation, slide_spec: Dict[str, Any], global_style: Dict[str, Any]) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide_style = _safe_get(slide_spec, "style", {})

    slide.shapes.title.text = _safe_get(slide_spec, "title", "")
    _set_title_text(slide.shapes.title, _safe_get(slide_spec, "title", ""), global_style, slide_style)

    if len(slide.placeholders) > 1:
        subtitle_text = _safe_get(slide_spec, "subtitle", "")
        slide.placeholders[1].text = subtitle_text
        _set_subtitle_text(slide.placeholders[1], subtitle_text, global_style, slide_style)

    if _safe_get(slide_style, "logo_required", True):
        _add_logo(slide, global_style, slide_style)

    _add_notes(slide, slide_spec)


def _add_bullet_slide(prs: Presentation, slide_spec: Dict[str, Any], global_style: Dict[str, Any]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide_style = _safe_get(slide_spec, "style", {})

    title_text = _safe_get(slide_spec, "title", "")
    slide.shapes.title.text = title_text
    _set_title_text(slide.shapes.title, title_text, global_style, slide_style)

    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    max_bullets = _style_value(global_style, slide_style, "max_bullets_per_slide", 5)
    bullets = _truncate_bullets(_safe_get(slide_spec, "bullets", []), max_bullets)

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        if isinstance(bullet, dict):
            p.text = str(bullet.get("text", ""))
            p.level = int(bullet.get("level", 0))
        else:
            p.text = str(bullet)
            p.level = 0

        _apply_paragraph_style(
            p,
            font_family=_style_value(global_style, slide_style, "body_font_family", "Arial"),
            font_size_pt=_style_value(global_style, slide_style, "body_font_size_pt", 18),
            bold=False,
            color_hex=_style_value(global_style, slide_style, "body_color", "#24135F"),
        )

    _add_notes(slide, slide_spec)


def _add_section_slide(prs: Presentation, slide_spec: Dict[str, Any], global_style: Dict[str, Any]) -> None:
    layout = prs.slide_layouts[2] if len(prs.slide_layouts) > 2 else prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide_style = _safe_get(slide_spec, "style", {})

    title_text = _safe_get(slide_spec, "title", "")
    slide.shapes.title.text = title_text
    _set_title_text(slide.shapes.title, title_text, global_style, slide_style)

    if len(slide.placeholders) > 1:
        subtitle_text = _safe_get(slide_spec, "subtitle", "")
        slide.placeholders[1].text = subtitle_text
        _set_subtitle_text(slide.placeholders[1], subtitle_text, global_style, slide_style)

    _add_notes(slide, slide_spec)


def _add_two_column_slide(prs: Presentation, slide_spec: Dict[str, Any], global_style: Dict[str, Any]) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    slide_style = _safe_get(slide_spec, "style", {})

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8.0), Inches(0.6))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = _safe_get(slide_spec, "title", "")
    _apply_paragraph_style(
        p_title,
        font_family=_style_value(global_style, slide_style, "title_font_family", "Arial"),
        font_size_pt=_style_value(global_style, slide_style, "title_font_size_pt", 28),
        bold=True,
        color_hex=_style_value(global_style, slide_style, "title_color", _style_value(global_style, slide_style, "primary_color", "#24135F")),
    )

    _add_logo(slide, global_style, slide_style)

    left_title = _safe_get(slide_spec, "left_title", "")
    right_title = _safe_get(slide_spec, "right_title", "")
    max_bullets = _style_value(global_style, slide_style, "max_bullets_per_slide", 5)
    left_bullets = _truncate_bullets(_safe_get(slide_spec, "left_bullets", []), max_bullets)
    right_bullets = _truncate_bullets(_safe_get(slide_spec, "right_bullets", []), max_bullets)

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(4.0), Inches(5.0))
    tf_left = left_box.text_frame
    p = tf_left.paragraphs[0]
    p.text = left_title
    _apply_paragraph_style(
        p,
        font_family=_style_value(global_style, slide_style, "body_font_family", "Arial"),
        font_size_pt=18,
        bold=True,
        color_hex=_style_value(global_style, slide_style, "accent_color", "#D0006F"),
    )

    for bullet in left_bullets:
        bp = tf_left.add_paragraph()
        bp.text = str(bullet)
        bp.level = 0
        _apply_paragraph_style(
            bp,
            font_family=_style_value(global_style, slide_style, "body_font_family", "Arial"),
            font_size_pt=_style_value(global_style, slide_style, "body_font_size_pt", 18),
            bold=False,
            color_hex=_style_value(global_style, slide_style, "body_color", "#24135F"),
        )

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.0), Inches(1.3), Inches(4.0), Inches(5.0))
    tf_right = right_box.text_frame
    p = tf_right.paragraphs[0]
    p.text = right_title
    _apply_paragraph_style(
        p,
        font_family=_style_value(global_style, slide_style, "body_font_family", "Arial"),
        font_size_pt=18,
        bold=True,
        color_hex=_style_value(global_style, slide_style, "accent_color", "#D0006F"),
    )

    for bullet in right_bullets:
        bp = tf_right.add_paragraph()
        bp.text = str(bullet)
        bp.level = 0
        _apply_paragraph_style(
            bp,
            font_family=_style_value(global_style, slide_style, "body_font_family", "Arial"),
            font_size_pt=_style_value(global_style, slide_style, "body_font_size_pt", 18),
            bold=False,
            color_hex=_style_value(global_style, slide_style, "body_color", "#24135F"),
        )

    _add_notes(slide, slide_spec)


def _add_image_slide(prs: Presentation, slide_spec: Dict[str, Any], global_style: Dict[str, Any]) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    slide_style = _safe_get(slide_spec, "style", {})

    title = _safe_get(slide_spec, "title", "")
    image_path = _safe_get(slide_spec, "image_path", "")
    caption = _safe_get(slide_spec, "caption", "")

    # If the image does not exist, degrade gracefully to a text slide instead of showing an ugly error
    if not image_path or not os.path.exists(image_path):
        fallback_spec = {
            "title": title or "Visual Placeholder",
            "bullets": [
                "No approved image asset was available at generation time.",
                "Replace this placeholder with a compliant visual if needed."
            ],
            "notes": _safe_get(slide_spec, "notes", "")
        }
        _add_bullet_slide(prs, fallback_spec, global_style)
        return

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8.0), Inches(0.6))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    _apply_paragraph_style(
        p_title,
        font_family=_style_value(global_style, slide_style, "title_font_family", "Arial"),
        font_size_pt=_style_value(global_style, slide_style, "title_font_size_pt", 28),
        bold=True,
        color_hex=_style_value(global_style, slide_style, "title_color", _style_value(global_style, slide_style, "primary_color", "#24135F")),
    )

    _add_logo(slide, global_style, slide_style)

    slide.shapes.add_picture(image_path, Inches(1), Inches(1.2), width=Inches(7.8))

    if caption:
        caption_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(7.8), Inches(0.4))
        tf_caption = caption_box.text_frame
        p_caption = tf_caption.paragraphs[0]
        p_caption.text = caption
        p_caption.alignment = PP_ALIGN.CENTER
        _apply_paragraph_style(
            p_caption,
            font_family=_style_value(global_style, slide_style, "body_font_family", "Arial"),
            font_size_pt=_style_value(global_style, slide_style, "caption_font_size_pt", 12),
            bold=False,
            color_hex=_style_value(global_style, slide_style, "body_color", "#24135F"),
        )

    _add_notes(slide, slide_spec)


def _build_presentation_from_spec(spec: Dict[str, Any]) -> Presentation:
    template_path = _safe_get(spec, "template_path", "")
    global_style = _safe_get(spec, "global_style", {})

    prs = Presentation(template_path) if template_path and os.path.exists(template_path) else Presentation()

    slides = _safe_get(spec, "slides", [])
    for slide_spec in slides:
        slide_type = _safe_get(slide_spec, "type", "bullet")

        if slide_type == "title":
            _add_title_slide(prs, slide_spec, global_style)
        elif slide_type == "section":
            _add_section_slide(prs, slide_spec, global_style)
        elif slide_type == "two_column":
            _add_two_column_slide(prs, slide_spec, global_style)
        elif slide_type == "image":
            _add_image_slide(prs, slide_spec, global_style)
        else:
            _add_bullet_slide(prs, slide_spec, global_style)

    return prs


@tool
def generate_presentation(spec_json: str) -> bytes:
    """
    Generate a PowerPoint presentation from a JSON specification.

    Expected top-level structure:
    {
      "template_path": "optional/path/to/template.pptx",
      "global_style": {
        "title_font_family": "Arial",
        "body_font_family": "Arial",
        "fallback_font_family": "Arial Unicode MS",
        "title_font_size_pt": 28,
        "body_font_size_pt": 18,
        "subtitle_font_size_pt": 20,
        "caption_font_size_pt": 12,
        "primary_color": "#24135F",
        "accent_color": "#D0006F",
        "body_color": "#24135F",
        "logo_file": "assets/logo_primary.png",
        "logo_position": "top_right",
        "max_bullets_per_slide": 5
      },
      "slides": [...]
    }

    Returns:
        bytes: PPTX presentation as bytes.
    """
    spec = json.loads(spec_json)
    prs = _build_presentation_from_spec(spec)

    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    return pptx_bytes.getvalue()