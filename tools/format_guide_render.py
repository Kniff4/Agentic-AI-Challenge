"""Format-guide presentation renderer.

Extracts content from a source PPTX and rebuilds it following layout/format rules
from a brand guidelines PDF. Preserves all original text while applying new structure.

Run locally:
    python tools/format_guide_render.py <source.pptx> <guidelines.pdf> <output.pptx>

Returns a new PPTX with source content restructured per guidelines.
"""
from ibm_watsonx_orchestrate.agent_builder.tools import tool, WXOFile

import argparse
import io
import json
import re
import xml.etree.ElementTree as ET
from pathlib import Path


# ─── XML namespaces ──────────────────────────────────────────────────────

NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}


# ─── Content extraction from source PPTX ─────────────────────────────────

def _extract_text_from_shape(shape):
    """Extract plain text from a shape's text frame."""
    try:
        if not hasattr(shape, 'text_frame'):
            return ""
        text_frame = shape.text_frame
        lines = []
        for paragraph in text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs).strip()
            if text:
                lines.append(text)
        return lines
    except Exception:
        return []


def _extract_slide_content(slide):
    """Extract structured content from a single slide.

    Returns: {
        "title": "...",
        "subtitle": "...",
        "bullets": ["...", "..."],
        "text_blocks": ["..."],  # free-form text boxes
    }
    """
    title = ""
    subtitle = ""
    bullets = []
    text_blocks = []

    shapes = list(slide.shapes)

    for i, shape in enumerate(shapes):
        lines = _extract_text_from_shape(shape)
        if not lines:
            continue

        # First text shape is likely title
        if i == 0 and not title:
            title = lines[0] if lines else ""
        # Second text shape might be subtitle
        elif i == 1 and not subtitle and lines:
            subtitle = lines[0]
        # Multi-line shapes are likely bullet lists
        elif len(lines) > 1:
            bullets.extend(lines)
        # Single-line non-title shapes are text blocks
        elif lines and lines[0] != title:
            text_blocks.extend(lines)

    return {
        "title": title,
        "subtitle": subtitle,
        "bullets": bullets,
        "text_blocks": text_blocks,
    }


def _extract_all_content(pptx_bytes):
    """Extract content from all slides in a PPTX.

    Returns: [{ "title": "...", "subtitle": "...", "bullets": [...] }, ...]
    """
    from pptx import Presentation

    prs = Presentation(io.BytesIO(pptx_bytes))
    slides = []

    for slide in prs.slides:
        content = _extract_slide_content(slide)
        if content.get("title") or content.get("bullets") or content.get("text_blocks"):
            slides.append(content)

    return slides


# ─── Guidelines PDF parsing ──────────────────────────────────────────────

def _parse_guidelines_pdf(pdf_bytes):
    """Parse brand guidelines PDF to extract layout rules and style specs.

    Returns: {
        "layouts": [
            { "name": "title", "description": "...", "max_bullets": 5 },
            ...
        ],
        "colors": { "primary": "#...", ... },
        "fonts": { "title": "Arial", ... },
        "spacing": { "margin_top": 0.5, ... },
    }
    """
    import pdfplumber

    text_content = ""
    try:
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            for page in pdf.pages:
                text_content += page.extract_text() or ""
    except Exception as e:
        raise ValueError(f"Failed to parse PDF: {e}")

    # Extract color specifications (look for hex or RGB patterns)
    colors = {}
    hex_matches = re.findall(r'#([0-9A-Fa-f]{6})', text_content)
    if hex_matches:
        colors['primary'] = f"#{hex_matches[0]}"
        if len(hex_matches) > 1:
            colors['accent'] = f"#{hex_matches[1]}"

    # Extract font names (look for common patterns)
    fonts = {}
    font_patterns = re.findall(r'\b(Arial|Helvetica|Times|Calibri|Verdana|Georgia|Roboto|Ubuntu)\b', text_content, re.I)
    if font_patterns:
        fonts['body'] = font_patterns[0]
        fonts['title'] = font_patterns[0]

    # Infer layout patterns from PDF text
    layouts = _infer_layouts_from_text(text_content)

    return {
        "layouts": layouts,
        "colors": colors,
        "fonts": fonts,
        "spacing": {},  # Can be enhanced with more sophisticated parsing
        "raw_text": text_content[:1000],  # Store excerpt for debugging
    }


def _infer_layouts_from_text(text):
    """Infer slide layout patterns from guideline text.

    Looks for keywords like "title", "content", "column", "layout", etc.
    """
    layouts = []

    # Common patterns to search for
    if re.search(r'title.*slide|slide.*title', text, re.I):
        layouts.append({
            "name": "title",
            "description": "Title slide with main title and subtitle",
            "max_bullets": 1,
        })

    if re.search(r'content.*slide|bullet|list', text, re.I):
        layouts.append({
            "name": "content",
            "description": "Content slide with title and bulleted list",
            "max_bullets": 5,
        })

    if re.search(r'two.*column|column|compare|side.?by.?side', text, re.I):
        layouts.append({
            "name": "two_column",
            "description": "Two-column comparison layout",
            "max_bullets": 3,
        })

    if re.search(r'closing|conclusion|thank|end', text, re.I):
        layouts.append({
            "name": "closing",
            "description": "Closing slide",
            "max_bullets": 3,
        })

    # If no patterns found, use sensible defaults
    if not layouts:
        layouts = [
            {"name": "title", "description": "Title slide", "max_bullets": 1},
            {"name": "content", "description": "Content slide", "max_bullets": 5},
            {"name": "closing", "description": "Closing slide", "max_bullets": 3},
        ]

    return layouts


# ─── Content-to-layout mapping ───────────────────────────────────────────

def _map_content_to_slides(extracted_content, guidelines):
    """Map extracted content to guideline layouts, splitting if necessary.

    If content exceeds layout capacity, splits across multiple slides of same type.

    Returns: [{ "type": "title", "title": "...", "bullets": [...] }, ...]
    """
    if not extracted_content:
        return []

    layouts = guidelines.get("layouts", [])
    if not layouts:
        layouts = [
            {"name": "content", "max_bullets": 5},
        ]

    output_slides = []
    is_first = True

    for content_slide in extracted_content:
        title = content_slide.get("title", "")
        bullets = content_slide.get("bullets", [])
        text_blocks = content_slide.get("text_blocks", [])
        subtitle = content_slide.get("subtitle", "")

        # First slide is title if it has subtitle or is clearly a title
        if is_first and (subtitle or not bullets):
            slide_type = "title"
            is_first = False
            max_bullets = _get_max_bullets_for_type(slide_type, layouts)
            output_slides.append({
                "type": slide_type,
                "title": title,
                "subtitle": subtitle or (bullets[0] if bullets else ""),
            })
            bullets = bullets[1:] if bullets else []
        else:
            is_first = False

        # Combine bullets and text blocks, split if needed
        all_content = bullets + text_blocks

        while all_content:
            slide_type = "content"
            if output_slides and not output_slides[-1].get("processed"):
                slide_type = "content"
            max_bullets = _get_max_bullets_for_type(slide_type, layouts)

            # Split content to fit layout capacity
            chunk = all_content[:max_bullets]
            all_content = all_content[max_bullets:]

            output_slides.append({
                "type": slide_type,
                "title": title if not output_slides else f"{title} (cont'd)",
                "bullets": chunk,
                "processed": True,
            })

        # Last slide is closing
        if output_slides and len(extracted_content) > 0 and content_slide == extracted_content[-1]:
            if output_slides[-1]["type"] != "closing":
                output_slides.append({
                    "type": "closing",
                    "title": "Thank you",
                    "bullets": [],
                })

    return output_slides


def _get_max_bullets_for_type(slide_type, layouts):
    """Get maximum bullet count for a slide type from guidelines."""
    for layout in layouts:
        if layout.get("name") == slide_type:
            return layout.get("max_bullets", 5)
    return 5


# ─── Rendering ──────────────────────────────────────────────────────────

def _render_slides(slides_spec, guidelines):
    """Render slides to PPTX using guidelines styling.

    Returns PPTX bytes.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Extract styling from guidelines
    colors = guidelines.get("colors", {})
    fonts = guidelines.get("fonts", {})

    primary_color = _hex_to_rgb(colors.get("primary", "#000000"))
    accent_color = _hex_to_rgb(colors.get("accent", "#0066CC"))
    title_font = fonts.get("title", "Arial")
    body_font = fonts.get("body", "Arial")

    for slide_spec in slides_spec:
        slide_type = slide_spec.get("type", "content")
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        if slide_type == "title":
            _render_title_slide(slide, slide_spec, primary_color, title_font, body_font)
        elif slide_type == "two_column":
            _render_two_column_slide(slide, slide_spec, primary_color, accent_color, title_font, body_font)
        else:  # content, closing
            _render_content_slide(slide, slide_spec, primary_color, accent_color, title_font, body_font)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _hex_to_rgb(hex_color):
    """Convert hex color to RGBColor."""
    from pptx.dml.color import RGBColor

    hex_color = hex_color.lstrip('#')
    if len(hex_color) != 6:
        hex_color = "000000"
    try:
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
    except ValueError:
        return RGBColor(0, 0, 0)


def _render_title_slide(slide, spec, primary_color, title_font, body_font):
    """Render a title slide."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.text = spec.get("title", "")
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    title_para.font.name = title_font

    # Subtitle
    if spec.get("subtitle"):
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = spec.get("subtitle", "")
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = RGBColor(64, 64, 64)
        subtitle_para.font.name = body_font


def _render_content_slide(slide, spec, primary_color, accent_color, title_font, body_font):
    """Render a content/bullet slide."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = spec.get("title", "")
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    title_para.font.name = title_font

    # Content/bullets
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    bullets = spec.get("bullets", [])
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = str(bullet)
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(32, 32, 32)
        p.font.name = body_font
        p.level = 0


def _render_two_column_slide(slide, spec, primary_color, accent_color, title_font, body_font):
    """Render a two-column comparison slide."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = spec.get("title", "")
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    title_para.font.name = title_font

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    left_bullets = spec.get("bullets", [])[:3]
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = str(bullet)
        p.font.size = Pt(16)
        p.font.name = body_font

    # Right column
    right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.5), Inches(5.9), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    right_bullets = spec.get("bullets", [])[3:6] if len(spec.get("bullets", [])) > 3 else []
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = str(bullet)
        p.font.size = Pt(16)
        p.font.name = body_font


# ─── Main orchestration ─────────────────────────────────────────────────

def format_guide_render_internal(source_pptx_bytes, guidelines_pdf_bytes):
    """Core formatting logic.

    Args:
        source_pptx_bytes: Source presentation content
        guidelines_pdf_bytes: Brand guidelines PDF

    Returns: Reformatted PPTX bytes
    """

    # Extract content from source
    content = _extract_all_content(source_pptx_bytes)

    # Parse guidelines
    guidelines = _parse_guidelines_pdf(guidelines_pdf_bytes)

    # Map content to layouts
    slides_spec = _map_content_to_slides(content, guidelines)

    # Render to PPTX
    output_bytes = _render_slides(slides_spec, guidelines)

    return output_bytes


# ─── Tool wrapper ───────────────────────────────────────────────────────

@tool(
    name="format_guide_render",
    description=(
        "Reformats a presentation to follow brand guidelines. "
        "REQUIRES: Two actual file uploads (not just names). "
        "Returns the reformatted .pptx as a downloadable file."
    ),
)
def format_guide_render(source_pptx: bytes, guidelines_pdf: bytes) -> dict:
    """Reformat presentation per brand guidelines."""
    try:
        # Check if the LLM mistakenly passed filenames as strings
        if isinstance(source_pptx, str) or isinstance(guidelines_pdf, str):
            return {"error": "You passed filenames instead of actual file objects. Please use the upload widgets to provide the files."}

        # Basic validation to ensure we have a PDF before parsing
        if not guidelines_pdf.startswith(b"%PDF"):
             return {"error": "The guidelines file is not a valid PDF. Please re-upload the brand guidelines .pdf file."}

        # Run your internal logic
        output_bytes = format_guide_render_internal(source_pptx, guidelines_pdf)

        # Returning a DICT with 'file_content' tells Orchestrate to show a download button
        return {
            "file_content": output_bytes,
            "file_name": "reformatted_presentation.pptx",
            "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        }
    except Exception as e:
        # Return the error to the chat so the LLM can explain it to the user
        return {"error": f"Processing failed: {str(e)}"}


# ─── CLI for local testing ──────────────────────────────────────────────

if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Reformat PPTX per brand guidelines PDF"
    )
    parser.add_argument("source_pptx", help="Path to source presentation")
    parser.add_argument("guidelines_pdf", help="Path to brand guidelines PDF")
    parser.add_argument("output_pptx", help="Where to write reformatted .pptx")
    args = parser.parse_args()

    source_bytes = Path(args.source_pptx).read_bytes()
    guidelines_bytes = Path(args.guidelines_pdf).read_bytes()

    print(guidelines_bytes[:1000])  # Debug: show PDF text excerpt

    output_bytes = format_guide_render_internal(source_bytes, guidelines_bytes)

    out_path = Path(args.output_pptx)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_bytes(output_bytes)
    print(f"Wrote {out_path} ({len(output_bytes):,} bytes)")
